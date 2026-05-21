from __future__ import annotations

import hashlib
import sys
from pathlib import Path
from urllib.request import urlretrieve


BASE_DIR = Path(__file__).resolve().parent
for candidate in [
    Path("/Users/wale/Desktop/softtttt/aiapis/apiweb/.venv/lib/python3.13/site-packages"),
    Path("/Users/wale/Desktop/softtttt/foodup/foodweb/.venv/lib/python3.13/site-packages"),
]:
    if candidate.exists() and str(candidate) not in sys.path:
        sys.path.append(str(candidate))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = BASE_DIR / "HEALTHCARE_LANGUAGE_PITCH_DECK.pptx"
ASSET_DIR = BASE_DIR / ".deck_assets"
REFERENCE_ASSET_DIR = Path("/Users/wale/Desktop/health/hweb/.deck_assets")


def rgb(code: str) -> RGBColor:
    code = code.replace("#", "")
    return RGBColor(int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))


PAPER = rgb("#EEF5F4")
WHITE = rgb("#FFFFFF")
INK = rgb("#0D2332")
SOFT = rgb("#607284")
LINE = rgb("#D4E1E8")
NAVY = rgb("#0F2E46")
COBALT = rgb("#1F5F93")
TEAL = rgb("#0C8A7B")
DEEP = rgb("#09151F")
DEEP_SOFT = rgb("#163149")
GOLD = rgb("#F1CF63")
MIST = rgb("#D8E6EE")


HERO_IMAGE = "https://images.unsplash.com/photo-1576091160550-2173dba999ef?auto=format&fit=crop&w=1600&q=80"
PATIENT_IMAGE = "https://images.unsplash.com/photo-1579684385127-1ef15d508118?auto=format&fit=crop&w=1400&q=80"
TEAM_IMAGE = "https://images.unsplash.com/photo-1576091160399-112ba8d25d1d?auto=format&fit=crop&w=1400&q=80"
DASHBOARD_IMAGE = "https://images.unsplash.com/photo-1579154204601-01588f351e67?auto=format&fit=crop&w=1400&q=80"


SLIDES = [
    {
        "type": "cover",
        "title": "Healthcare Language",
        "subtitle": "A cloud operating system for multilingual care communication, translation workflows, and governed AI language support.",
        "eyebrow": "Language + Healthcare + AI + Data + Cloud",
        "image": HERO_IMAGE,
        "stats": [
            ("$288.6B", "digital health"),
            ("$866.5B", "healthcare IT"),
            ("$42.5B", "health cloud"),
        ],
    },
    {
        "type": "three_panel",
        "title": "Problem",
        "subtitle": "Healthcare systems still manage language access, translated content, and patient communication across fragmented tools and manual workflows.",
        "items": [
            ("Language access is disconnected", "Interpreter requests, translated instructions, reminders, and patient outreach often sit outside core care and administrative workflows."),
            ("Quality and terminology drift", "Hospitals and programs struggle to keep multilingual content, approved phrases, and clinical terminology aligned across teams and vendors."),
            ("Communication data is underused", "Language demand, outreach performance, and patient-understanding signals rarely live inside one governed operational system."),
        ],
    },
    {
        "type": "signal",
        "title": "Solution",
        "subtitle": "Healthcare Language gives providers and governments one secure layer for multilingual communication, content governance, and assistive language AI.",
        "image": PATIENT_IMAGE,
        "bullets": [
            "Patient language profiles, communication preferences, and access needs in one system",
            "Clinical translation, review, approval, and publishing workflows for regulated content",
            "Speech, transcription, summarization, and translation assistance with human oversight",
            "Dashboards for language demand, service gaps, quality control, and program performance",
        ],
    },
    {
        "type": "three_panel",
        "title": "Why now",
        "subtitle": "Three market forces make healthcare language infrastructure a timely, durable modernization opportunity.",
        "items": [
            ("Digital health is scaling", "The global digital health market reached $288.6B in 2024 and is projected to reach $946.0B by 2030."),
            ("Healthcare IT budgets are large", "The global healthcare IT market is estimated at $866.5B in 2025 and projected to exceed $2.86T by 2033."),
            ("AI + cloud need governance", "Healthcare cloud reached $42.5B in 2024 while AI in healthcare is projected to exceed $505.6B by 2033."),
        ],
    },
    {
        "type": "map",
        "title": "Product wedge",
        "subtitle": "Healthcare Language sits at the center of patient communication, translation operations, and multilingual service oversight.",
        "left_title": "Inputs",
        "left_points": ["Patient language data", "Clinical content", "Outreach demand"],
        "center_title": "Healthcare Language core",
        "center_points": ["Profiles", "Translation", "Terminology", "AI support", "Oversight"],
        "right_title": "Outputs",
        "right_points": ["Better understanding", "Faster delivery", "Quality control", "Access analytics"],
    },
    {
        "type": "stack",
        "title": "Core modules",
        "subtitle": "The platform expands from a communication wedge into broader multilingual health infrastructure over time.",
        "layers": [
            ("Patient access layer", "Preferred language, interpreter requirements, communication preferences, and accessibility context."),
            ("Translation workflow layer", "Request intake, reviewer routing, version control, approvals, and regulated content publishing."),
            ("Terminology layer", "Medical glossaries, approved phrases, reusable content, and multilingual quality governance."),
            ("Intelligence layer", "Demand forecasting, quality review, audit trails, summarization, and governed AI assistance."),
        ],
    },
    {
        "type": "three_panel",
        "title": "Beachhead customers",
        "subtitle": "The first buyers are healthcare organizations with high communication complexity, compliance needs, and population-scale service delivery.",
        "items": [
            ("Public health programs", "Programs need multilingual outreach, campaign control, reviewed patient content, and service continuity across diverse populations."),
            ("Hospital and insurer networks", "Institutions need one governed environment for patient communication, terminology, translation workflows, and access reporting."),
            ("Residential and long-term care", "Care operators need clearer instructions, caregiver coordination, and multilingual communication support across settings."),
        ],
    },
    {
        "type": "timeline",
        "title": "How value is created",
        "subtitle": "The product creates repeatable operational value from language intake to communication quality and system insight.",
        "steps": [
            ("Capture", "Record patient language needs, access preferences, and communication context."),
            ("Coordinate", "Route translation requests, interpreter support, content review, and care messaging workflows."),
            ("Govern", "Standardize terminology, approvals, publishing, and auditability across programs and institutions."),
            ("Improve", "Use analytics and AI support to forecast demand, flag quality issues, and strengthen access outcomes."),
        ],
    },
    {
        "type": "ledger",
        "title": "Why this can win",
        "subtitle": "Workflow depth, language governance, and regulated deployment needs create defensibility once embedded in care operations.",
        "rows": [
            ("Access control", "Patient communication and language records", "defensible"),
            ("Audit trail", "Translation edits and approvals", "embedded"),
            ("Terminology governance", "Approved medical language and content", "sticky"),
            ("Operational reporting", "Demand, quality, and access monitoring", "required"),
        ],
    },
    {
        "type": "three_panel",
        "title": "Market opportunity",
        "subtitle": "Healthcare Language sits inside multiple large modernization markets with recurring demand for better communication infrastructure.",
        "items": [
            ("$288.6B digital health", "Global digital health reached $288.6B in 2024 and is projected to grow to $946.0B by 2030."),
            ("$866.5B healthcare IT", "Healthcare IT is estimated at $866.5B in 2025, with projected revenue above $2.86T by 2033."),
            ("$505.6B AI in healthcare", "AI in healthcare is projected to exceed $505.6B by 2033 as workflows and analytics move into governed systems."),
        ],
    },
    {
        "type": "signal",
        "title": "Go-to-market",
        "subtitle": "Commercial entry starts with urgent language-access and communication workflows, then expands into broader system integrations.",
        "image": TEAM_IMAGE,
        "bullets": [
            "Lead with multilingual patient communication, translation workflow, and terminology control use cases",
            "Sell into public-sector programs, hospital groups, insurers, and regulated care operators",
            "Expand into outreach analytics, speech workflows, AI assistance, and network-wide governance",
            "Build recurring revenue through deployment, configuration, integrations, and enterprise contracts",
        ],
    },
    {
        "type": "three_panel",
        "title": "Revenue model",
        "subtitle": "Revenue compounds through platform contracts, implementation services, and high-value language intelligence modules.",
        "items": [
            ("Platform licensing", "Contracts tied to institution, network, program, or regional deployment for multilingual communication infrastructure."),
            ("Implementation revenue", "Integration, content migration, rollout, terminology setup, training, and workflow configuration services."),
            ("Expansion modules", "Premium analytics, speech tools, reviewed AI assistance, outreach optimization, and multilingual quality monitoring."),
        ],
    },
    {
        "type": "timeline",
        "title": "Expansion roadmap",
        "subtitle": "The business expands from language access workflows into broader communication infrastructure and operational intelligence.",
        "steps": [
            ("Phase 1", "Win with translation workflows, patient language profiles, and governed content operations."),
            ("Phase 2", "Expand into messaging, outreach, interpreter coordination, and network-wide reporting."),
            ("Phase 3", "Layer in transcription, summarization, quality review, and assistive AI workflows."),
            ("Phase 4", "Scale to system-wide communication intelligence, optimization, and broader care integration."),
        ],
    },
    {
        "type": "signal",
        "title": "Investment case",
        "subtitle": "Healthcare Language sits where digital health modernization, communication equity, and governed healthcare AI converge.",
        "image": DASHBOARD_IMAGE,
        "bullets": [
            "Large multi-market opportunity across digital health, healthcare IT, AI, and cloud infrastructure",
            "Workflow-embedded product with strong stickiness once language operations move into the platform",
            "Recurring enterprise revenue with implementation, integrations, and expansion layers",
            "Clear fit for public-sector, provider, payer, and regulated care environments",
        ],
    },
    {
        "type": "close",
        "title": "Healthcare Language",
        "subtitle": "We are building the infrastructure layer for multilingual healthcare communication, translation operations, and governed language AI.",
        "contact": "enquiry@healthcarelanguage.ca",
        "image": HERO_IMAGE,
    },
]


def add_textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def ensure_image(url: str) -> Path:
    ASSET_DIR.mkdir(exist_ok=True)
    digest = hashlib.sha1(url.encode("utf-8")).hexdigest()[:16]
    path = ASSET_DIR / f"{digest}.jpg"
    if path.exists():
        return path
    reference_path = REFERENCE_ASSET_DIR / path.name
    if reference_path.exists():
        path.write_bytes(reference_path.read_bytes())
        return path
    urlretrieve(url, path)
    return path


def add_slide_number(slide, number: int):
    add_textbox(slide, Inches(10.25), Inches(5.9), Inches(0.45), Inches(0.18), f"{number:02d}", size=9, color=SOFT, font="Aptos", align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle, light=False):
    accent = GOLD if not light else TEAL
    title_color = WHITE if light else INK
    sub_color = MIST if light else SOFT
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(0.72), Inches(0.14), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, Inches(0.93), Inches(0.72), Inches(9.8), Inches(0.45), title, size=27, color=title_color, bold=True, font="Georgia")
    add_textbox(slide, Inches(0.93), Inches(1.23), Inches(9.7), Inches(0.56), subtitle, size=12, color=sub_color, font="Aptos")


def cover_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_card(slide, Inches(0.42), Inches(0.42), Inches(10.48), Inches(5.55), DEEP)
    add_card(slide, Inches(8.08), Inches(0.42), Inches(2.82), Inches(5.55), DEEP_SOFT)
    if data.get("image"):
        slide.shapes.add_picture(str(ensure_image(data["image"])), Inches(8.08), Inches(0.42), width=Inches(2.82), height=Inches(5.55))
    star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, Inches(8.72), Inches(0.95), Inches(1.48), Inches(1.48))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.color.rgb = GOLD
    add_textbox(slide, Inches(0.9), Inches(0.9), Inches(3.6), Inches(0.3), data["eyebrow"], size=13, color=rgb("#B8D3E6"), bold=True, font="Aptos")
    add_textbox(slide, Inches(0.88), Inches(1.52), Inches(6.4), Inches(0.95), data["title"], size=31, color=WHITE, bold=True, font="Georgia")
    add_textbox(slide, Inches(0.9), Inches(2.72), Inches(6.15), Inches(1.0), data["subtitle"], size=16, color=MIST, font="Aptos")
    add_textbox(slide, Inches(0.92), Inches(5.0), Inches(4.2), Inches(0.22), "enquiry@healthcarelanguage.ca", size=12, color=MIST, bold=True, font="Aptos")
    for idx, (label, copy) in enumerate(data["stats"]):
        left = Inches(8.32)
        top = Inches(2.0 + idx * 1.16)
        add_card(slide, left, top, Inches(2.18), Inches(0.92), rgb("#1B3750"))
        add_textbox(slide, left + Inches(0.18), top + Inches(0.14), Inches(1.7), Inches(0.2), label, size=10, color=rgb("#A6C8DF"), bold=True, font="Aptos")
        add_textbox(slide, left + Inches(0.18), top + Inches(0.4), Inches(1.8), Inches(0.26), copy, size=12, color=WHITE, bold=True, font="Georgia")
    add_slide_number(slide, number)


def three_panel_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    for idx, (heading, copy) in enumerate(data["items"]):
        left = Inches(0.65 + idx * 3.55)
        add_card(slide, left, Inches(2.08), Inches(3.2), Inches(3.75), WHITE, LINE)
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, Inches(2.08), Inches(3.2), Inches(0.12))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD if idx != 1 else TEAL
        accent.line.color.rgb = GOLD if idx != 1 else TEAL
        add_textbox(slide, left + Inches(0.22), Inches(2.33), Inches(2.7), Inches(0.36), heading, size=18, color=NAVY, bold=True, font="Georgia")
        add_textbox(slide, left + Inches(0.22), Inches(2.92), Inches(2.75), Inches(1.52), copy, size=12, color=SOFT, font="Aptos")
    add_slide_number(slide, number)


def signal_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_card(slide, Inches(0.52), Inches(0.52), Inches(10.3), Inches(5.9), PAPER)
    add_title(slide, data["title"], data["subtitle"])
    has_image = bool(data.get("image"))
    if has_image:
        add_card(slide, Inches(6.55), Inches(1.92), Inches(3.65), Inches(3.85), WHITE, LINE)
        slide.shapes.add_picture(str(ensure_image(data["image"])), Inches(6.72), Inches(2.08), width=Inches(3.31), height=Inches(3.52))
    for idx, bullet in enumerate(data["bullets"]):
        y = Inches(2.16 + idx * 0.92)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95), y + Inches(0.08), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.color.rgb = GOLD
        add_textbox(slide, Inches(1.28), y, Inches(4.95 if has_image else 8.7), Inches(0.42), bullet, size=16, color=INK, font="Aptos")
    add_slide_number(slide, number)


def map_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    columns = [
        (Inches(0.85), Inches(2.08), Inches(2.55), Inches(3.55), WHITE, data["left_title"], data["left_points"]),
        (Inches(4.1), Inches(1.82), Inches(3.25), Inches(4.05), NAVY, data["center_title"], data["center_points"]),
        (Inches(8.05), Inches(2.08), Inches(2.55), Inches(3.55), WHITE, data["right_title"], data["right_points"]),
    ]
    for left, top, width, height, fill, title, points in columns:
        add_card(slide, left, top, width, height, fill, LINE if fill == WHITE else NAVY)
        title_color = WHITE if fill == NAVY else INK
        body_color = MIST if fill == NAVY else SOFT
        add_textbox(slide, left + Inches(0.22), top + Inches(0.2), width - Inches(0.44), Inches(0.35), title, size=18, color=title_color, bold=True, font="Georgia")
        for idx, point in enumerate(points):
            add_textbox(slide, left + Inches(0.28), top + Inches(0.84 + idx * 0.58), width - Inches(0.56), Inches(0.24), f"• {point}", size=12, color=body_color, font="Aptos")
    add_slide_number(slide, number)


def stack_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, data["title"], data["subtitle"])
    fills = [rgb("#EAF2F6"), rgb("#DDEBF3"), rgb("#EAF2F6"), rgb("#DDEBF3")]
    for idx, (heading, copy) in enumerate(data["layers"]):
        top = Inches(2.0 + idx * 1.0)
        add_card(slide, Inches(0.78), top, Inches(9.95), Inches(0.78), fills[idx], rgb("#C9DDE7"))
        add_textbox(slide, Inches(1.02), top + Inches(0.16), Inches(2.4), Inches(0.22), heading, size=15, color=NAVY, bold=True, font="Georgia")
        add_textbox(slide, Inches(3.2), top + Inches(0.14), Inches(7.1), Inches(0.26), copy, size=11, color=INK, font="Aptos")
    add_slide_number(slide, number)


def timeline_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP
    add_title(slide, data["title"], data["subtitle"], light=True)
    for idx, (heading, copy) in enumerate(data["steps"]):
        left = Inches(0.82 + idx * 2.48)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.58), Inches(2.1), Inches(2.48))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb("#173044")
        shape.line.color.rgb = rgb("#22445F")
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), Inches(2.82), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GOLD
        badge.line.color.rgb = GOLD
        add_textbox(slide, left + Inches(0.62), Inches(2.82), Inches(0.4), Inches(0.18), f"{idx + 1:02d}", size=11, color=rgb("#A7C9DF"), bold=True, font="Aptos")
        add_textbox(slide, left + Inches(0.18), Inches(3.25), Inches(1.72), Inches(0.28), heading, size=17, color=WHITE, bold=True, font="Georgia")
        add_textbox(slide, left + Inches(0.18), Inches(3.8), Inches(1.72), Inches(0.78), copy, size=11, color=MIST, font="Aptos")
    add_slide_number(slide, number)


def ledger_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, data["title"], data["subtitle"])
    add_card(slide, Inches(0.72), Inches(2.0), Inches(10.0), Inches(3.82), WHITE, LINE)
    add_textbox(slide, Inches(1.0), Inches(2.24), Inches(2.3), Inches(0.18), "Control", size=11, color=SOFT, bold=True, font="Aptos")
    add_textbox(slide, Inches(4.4), Inches(2.24), Inches(2.8), Inches(0.18), "Record or workflow", size=11, color=SOFT, bold=True, font="Aptos")
    add_textbox(slide, Inches(8.92), Inches(2.24), Inches(1.1), Inches(0.18), "Status", size=11, color=SOFT, bold=True, font="Aptos")
    for idx, row in enumerate(data["rows"]):
        y = Inches(2.74 + idx * 0.72)
        rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.94), y - Inches(0.06), Inches(9.55), Inches(0.01))
        rule.fill.solid()
        rule.fill.fore_color.rgb = rgb("#D9E7EE")
        rule.line.color.rgb = rgb("#D9E7EE")
        add_textbox(slide, Inches(1.0), y, Inches(2.85), Inches(0.22), row[0], size=12, color=INK, font="Aptos")
        add_textbox(slide, Inches(4.4), y, Inches(3.1), Inches(0.22), row[1], size=12, color=SOFT, font="Aptos")
        add_textbox(slide, Inches(8.98), y, Inches(1.0), Inches(0.2), row[2], size=10, color=TEAL, bold=True, font="Aptos")
    add_slide_number(slide, number)


def close_slide(prs, data, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DEEP
    add_textbox(slide, Inches(0.9), Inches(1.08), Inches(3.4), Inches(0.26), "Healthcare Communication Infrastructure", size=14, color=rgb("#B8D3E6"), bold=True, font="Aptos")
    add_textbox(slide, Inches(0.88), Inches(1.88), Inches(5.8), Inches(0.8), data["title"], size=30, color=WHITE, bold=True, font="Georgia")
    add_textbox(slide, Inches(0.9), Inches(2.88), Inches(6.45), Inches(0.95), data["subtitle"], size=16, color=MIST, font="Aptos")
    add_textbox(slide, Inches(0.92), Inches(4.95), Inches(4.2), Inches(0.22), data["contact"], size=14, color=WHITE, bold=True, font="Aptos")
    add_card(slide, Inches(7.92), Inches(1.0), Inches(2.42), Inches(4.92), DEEP_SOFT)
    if data.get("image"):
        slide.shapes.add_picture(str(ensure_image(data["image"])), Inches(7.92), Inches(1.0), width=Inches(2.42), height=Inches(4.92))
    star = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.STAR_5_POINT, Inches(8.43), Inches(1.48), Inches(1.4), Inches(1.4))
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.color.rgb = GOLD
    add_textbox(slide, Inches(8.12), Inches(3.64), Inches(2.0), Inches(0.24), "Governed language", size=11, color=MIST, bold=True, font="Aptos", align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.12), Inches(4.0), Inches(2.0), Inches(0.48), "care, data,\nand AI", size=15, color=WHITE, bold=True, font="Georgia", align=PP_ALIGN.CENTER)
    add_slide_number(slide, number)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(11.33)
    prs.slide_height = Inches(6.38)

    for idx, slide_data in enumerate(SLIDES, start=1):
        kind = slide_data["type"]
        if kind == "cover":
            cover_slide(prs, slide_data, idx)
        elif kind == "three_panel":
            three_panel_slide(prs, slide_data, idx)
        elif kind == "signal":
            signal_slide(prs, slide_data, idx)
        elif kind == "map":
            map_slide(prs, slide_data, idx)
        elif kind == "stack":
            stack_slide(prs, slide_data, idx)
        elif kind == "timeline":
            timeline_slide(prs, slide_data, idx)
        elif kind == "ledger":
            ledger_slide(prs, slide_data, idx)
        elif kind == "close":
            close_slide(prs, slide_data, idx)

    prs.core_properties.title = "Healthcare Language Pitch Deck"
    prs.core_properties.subject = "Multilingual healthcare communication infrastructure"
    prs.core_properties.author = "OpenAI Codex"
    prs.save(OUT)


if __name__ == "__main__":
    build_deck()
    print(OUT)
